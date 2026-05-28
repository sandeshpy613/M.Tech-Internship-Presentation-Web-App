import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Define paths
output_dir = r"c:\Users\Sandesh PY\Desktop\Internship Portfolio\Project\public\presentation"
os.makedirs(output_dir, exist_ok=True)
ppt_path = os.path.join(output_dir, "internship_presentation.pptx")

# Load single source of truth slide-content.json
json_path = r"c:\Users\Sandesh PY\Desktop\Internship Portfolio\Project\components\presentation\slide-content.json"
with open(json_path, "r", encoding="utf-8") as f:
    data = json.load(f)

info = data["presentation_info"]
s_data = {slide["key"]: slide for slide in data["slides"]}

# Initialize presentation
prs = Presentation()
# Set 16:9 widescreen layout
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Styling Helper Variables - using Times New Roman for all text
font_title = "Times New Roman"
font_body = "Times New Roman"
color_title = RGBColor(15, 23, 42)      # Deep slate
color_body = RGBColor(71, 85, 105)     # Muted slate
color_accent = RGBColor(37, 99, 235)    # Blue

def apply_background(slide):
    # Set slide background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

# ----------------- SLIDE 1: Title Slide -----------------
slide_layout = prs.slide_layouts[6] # Blank layout
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)

# Add College Logo at top - centered and wide (11.33 inches wide, matching text boundaries)
logo_path = r"c:\Users\Sandesh PY\Desktop\Internship Portfolio\Project\header\logo.jpg"
slide.shapes.add_picture(logo_path, Inches(1.0), Inches(0.4), width=Inches(11.33))

# Main Title & Subtitle Box - pushed down to avoid overlapping the logo banner
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(4.5))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = info["title"]
p1.font.name = font_title
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = color_title
p1.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = f"{info['subtitle']}  |  {info['affiliation']}"
p2.font.name = font_body
p2.font.size = Pt(18)
p2.font.color.rgb = color_accent
p2.space_after = Pt(40)

p3 = tf.add_paragraph()
p3.text = f"Presenter: {info['presenter']['name']}  (Reg No: {info['presenter']['reg_no']})\nGuide: {info['guide']['name']} ({info['guide']['designation']})\nExternal Guide: {info['external_guide']['name']} ({info['external_guide']['company']})"
p3.font.name = font_body
p3.font.size = Pt(14)
p3.font.color.rgb = color_body

# ----------------- SLIDE Helper for Standard Slides -----------------
def add_standard_slide(title, category=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(0.8))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.name = font_title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = color_title
    
    # Draw a little bolded horizontal separator line under title (accent blue)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.65), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color_accent
    line.line.color.rgb = color_accent
    
    # Footer info
    footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.33), Inches(0.4))
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = f"{info['presenter']['name']}  •  M.TECH INTERNSHIP PRESENTATION"
    fp.font.name = font_body
    fp.font.size = Pt(9)
    fp.font.color.rgb = RGBColor(148, 163, 184)
    
    return slide

# ----------------- SLIDE 2: Agenda (Table of Contents) -----------------
slide = add_standard_slide("Agenda", "")
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Presentation Structure:"
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(20)

agenda_items = [
    "1. Executive Summary",
    "2. Company Context (Raashi Digital LLP)",
    "3. Internship Training Domains",
    "4. Theoretical Context & Practical Coding Skills",
    "5. Project 1: Enterprise Learning Portal (Frontend, Backend, Database)",
    "6. Project 2: AI-Driven Knowledge System (RAG Ingestion)",
    "7. Grounding and Search Priorities (AI Safety & Relevance)",
    "8. Conclusion & Future Outlook"
]

for item in agenda_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# ----------------- SLIDE 3: Executive Summary -----------------
slide = add_standard_slide(s_data["executive_summary"]["title"], s_data["executive_summary"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = s_data["executive_summary"]["heading"]
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(20)

for b in s_data["executive_summary"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.name = font_body
    p.font.size = Pt(16)
    p.font.color.rgb = color_body
    p.space_after = Pt(12)

# ----------------- SLIDE 4: Organizational Context -----------------
slide = add_standard_slide(s_data["company_context"]["title"], s_data["company_context"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = s_data["company_context"]["heading"]
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(20)

for b in s_data["company_context"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.name = font_body
    p.font.size = Pt(16)
    p.font.color.rgb = color_body
    p.space_after = Pt(12)

# ----------------- SLIDE 5: Training Program -----------------
slide = add_standard_slide(s_data["training_program"]["title"], s_data["training_program"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Domain 1: " + s_data["training_program"]["domain_1"]["title"]
p.font.name = font_body
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = color_accent
p.space_after = Pt(8)

for b in s_data["training_program"]["domain_1"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "  - " + b
    p.font.name = font_body
    p.font.size = Pt(14)
    p.font.color.rgb = color_body
    p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "Domain 2: " + s_data["training_program"]["domain_2"]["title"]
p.font.name = font_body
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = color_accent
p.space_before = Pt(18)
p.space_after = Pt(8)

for b in s_data["training_program"]["domain_2"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "  - " + b
    p.font.name = font_body
    p.font.size = Pt(14)
    p.font.color.rgb = color_body
    p.space_after = Pt(6)

# ----------------- SLIDE 6: Learning Experiences -----------------
slide = add_standard_slide(s_data["learning_experiences"]["title"], s_data["learning_experiences"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

th = s_data["learning_experiences"]["theory"]
p = tf.paragraphs[0]
p.text = th["heading"] + ":"
p.font.name = font_body
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(8)

for item in th["items"]:
    p = tf.add_paragraph()
    p.text = f"  • {item['title']}: {item['desc']}"
    p.font.name = font_body
    p.font.size = Pt(14)
    p.font.color.rgb = color_body
    p.space_after = Pt(6)

pr = s_data["learning_experiences"]["practical"]
p = tf.add_paragraph()
p.text = pr["heading"] + ":"
p.font.name = font_body
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = color_title
p.space_before = Pt(18)
p.space_after = Pt(8)

for item in pr["items"]:
    p = tf.add_paragraph()
    p.text = "  • " + item
    p.font.name = font_body
    p.font.size = Pt(14)
    p.font.color.rgb = color_body
    p.space_after = Pt(6)

# ----------------- SLIDE 7: Project 1 - Frontend -----------------
slide = add_standard_slide(s_data["project_1"]["frontend"]["title"], s_data["project_1"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = s_data["project_1"]["frontend"]["subtitle"] + ":"
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(16)

for b in s_data["project_1"]["frontend"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# ----------------- SLIDE 8: Project 1 - Backend -----------------
slide = add_standard_slide(s_data["project_1"]["backend"]["title"], s_data["project_1"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = s_data["project_1"]["backend"]["subtitle"] + ":"
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(16)

for b in s_data["project_1"]["backend"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# ----------------- SLIDE 9: Project 1 - Database -----------------
slide = add_standard_slide(s_data["project_1"]["database"]["title"], s_data["project_1"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = s_data["project_1"]["database"]["subtitle"] + ":"
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(16)

for b in s_data["project_1"]["database"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# ----------------- SLIDE 10: Project 2 - AI Knowledge System -----------------
slide = add_standard_slide(s_data["project_2"]["title"], s_data["project_2"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = s_data["project_2"]["heading"]
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(16)

for pt in s_data["project_2"]["bullets"]:
    p = tf.add_paragraph()
    p.text = "• " + pt
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# ----------------- SLIDE 11: Grounding and Search Priorities -----------------
slide = add_standard_slide(s_data["grounding"]["title"], s_data["grounding"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "How the System Enforces Safety and Priority:"
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(16)

points = [
    s_data["grounding"]["priority"]["heading"] + ": " + s_data["grounding"]["priority"]["desc"] + " " + s_data["grounding"]["priority"]["highlight"],
    "Cross-Encoder Re-Ranking: " + s_data["grounding"]["grounding"]["items"][0],
    "Strict Grounding: " + s_data["grounding"]["grounding"]["items"][1],
    "Source Citations: " + s_data["grounding"]["grounding"]["items"][2]
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "• " + pt
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# ----------------- SLIDE 12: Conclusion & Outlook -----------------
slide = add_standard_slide(s_data["conclusion"]["title"], s_data["conclusion"]["category"])
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Key Takeaways & Future Trajectory:"
p.font.name = font_body
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = color_title
p.space_after = Pt(16)

points = [
    f"{s_data['conclusion']['capstone']['title']} ({s_data['conclusion']['capstone']['subtitle']}): {s_data['conclusion']['capstone']['desc']}",
    f"{s_data['conclusion']['foundation']['title']} ({s_data['conclusion']['foundation']['subtitle']}): {s_data['conclusion']['foundation']['desc']}",
    f"{s_data['conclusion']['future']['title']} ({s_data['conclusion']['future']['subtitle']}): {s_data['conclusion']['future']['desc']}"
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "• " + pt
    p.font.name = font_body
    p.font.size = Pt(15)
    p.font.color.rgb = color_body
    p.space_after = Pt(10)

# Save presentation
prs.save(ppt_path)
print("PowerPoint presentation successfully created at:", ppt_path)
